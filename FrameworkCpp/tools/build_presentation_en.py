"""Generate the English UAM MDO presentation (mirrors UAM_MDO_Presentation.pptx style).

Content is grounded in the live FrameworkCpp source:
  - 16 active design variables  (HexacopterArchitecture::registerDefaultParameters)
  - 7 active weighted objective terms (EvaluationContext::objective_weights)
  - 18 architecture-level hard constraints (HexacopterArchitecture::registerDefaultConstraints)
Run:  python tools/build_presentation_en.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- palette (sampled from the Korean reference deck) -------------------------
NAVY    = RGBColor(0x0D, 0x1B, 0x3E)
NAVY2   = RGBColor(0x1A, 0x23, 0x40)
BLUE    = RGBColor(0x14, 0x51, 0xA8)
BLUE_LT = RGBColor(0x1A, 0x3A, 0x6E)
CYAN    = RGBColor(0x0B, 0x8F, 0xD4)
CYAN2   = RGBColor(0x00, 0xC4, 0xE8)
ORANGE  = RGBColor(0xF0, 0x7C, 0x26)
GREEN   = RGBColor(0x3F, 0xA0, 0x6B)
GRAY    = RGBColor(0x4A, 0x5A, 0x78)
GRAY_LT = RGBColor(0x7A, 0x8B, 0xAA)
FILL1   = RGBColor(0xED, 0xF3, 0xFB)
FILL2   = RGBColor(0xD8, 0xE4, 0xF0)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
INK     = RGBColor(0x20, 0x2A, 0x44)

FONT = "Calibri"
FONT_L = "Calibri Light"

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def rect(s, l, t, w, h, fill=None, line=None, lw=0.75, shape=MSO_SHAPE.RECTANGLE):
    sp = s.shapes.add_shape(shape, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(lw)
    return sp


def txt(s, l, t, w, h, runs, size=12, color=INK, bold=False, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, font=FONT, sp_after=2, line=1.0):
    """runs: str  OR list of paragraphs; each paragraph is str or list of (text,opts) run tuples."""
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    if isinstance(runs, str):
        runs = [runs]
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(sp_after)
        p.space_before = Pt(0)
        p.line_spacing = line
        chunks = para if isinstance(para, list) else [(para, {})]
        for text, o in chunks:
            r = p.add_run(); r.text = text
            r.font.name = o.get("font", font)
            r.font.size = Pt(o.get("size", size))
            r.font.bold = o.get("bold", bold)
            r.font.color.rgb = o.get("color", color)
    return tb


def header(s, title, num):
    rect(s, 0, 0, 10, 0.62, fill=NAVY)
    rect(s, 0, 0, 0.16, 0.62, fill=ORANGE)
    txt(s, 0.36, 0, 8.6, 0.62, title, size=19, color=WHITE, bold=True,
        anchor=MSO_ANCHOR.MIDDLE, font=FONT_L)
    txt(s, 9.0, 0, 0.85, 0.62, num, size=16, color=CYAN2, bold=True,
        anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)


def footer(s):
    txt(s, 0.36, 5.32, 9.3, 0.26, "Multidisciplinary Design Optimization of Aerospace Systems  ·  FrameworkCpp  ·  2026",
        size=8, color=GRAY_LT)


def card(s, l, t, w, h, accent, fill=FILL1):
    rect(s, l, t, w, h, fill=fill)
    rect(s, l, t, 0.07, h, fill=accent)


# =============================================================================
# SLIDE 1 — Title
# =============================================================================
s = slide()
rect(s, 0, 0, 10, 5.625, fill=NAVY)
rect(s, 0, 0, 10, 0.10, fill=ORANGE)
rect(s, 0, 5.525, 10, 0.10, fill=CYAN)
txt(s, 0.6, 1.05, 8.8, 1.0, "UAM Hexacopter", size=46, color=WHITE, bold=True, font=FONT_L)
txt(s, 0.6, 1.95, 8.8, 0.7, "Multidisciplinary Design Optimization Framework", size=26,
    color=CYAN2, bold=True, font=FONT_L)
rect(s, 0.62, 2.78, 2.4, 0.04, fill=ORANGE)
txt(s, 0.6, 2.95, 9.0, 0.4,
    "FrameworkCpp  ·  A C++ MDO System for eVTOL Hexacopter Design Optimization",
    size=13, color=FILL2)
# chips
chips = [("Pagmo2 · CMA-ES / NSGA-II", CYAN), ("ACS · Structural · Powertrain", BLUE),
         ("16 design variables", ORANGE), ("18 hard constraints", GREEN)]
x = 0.6
for label, c in chips:
    w = 0.18 + 0.092 * len(label)
    rect(s, x, 4.05, w, 0.36, fill=NAVY2, line=c, lw=1.25)
    txt(s, x, 4.05, w, 0.36, label, size=10.5, color=WHITE, align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE)
    x += w + 0.18
txt(s, 0.6, 4.95, 9.0, 0.3, "Multidisciplinary Design Optimization of Aerospace Systems  ·  2026",
    size=11, color=GRAY_LT)

# =============================================================================
# SLIDE 2 — Project Overview
# =============================================================================
s = slide(); header(s, "Project Overview", "02"); footer(s)
items = [
    ("①", "Study Object",
     ["UAM (Urban Air Mobility) hexacopter",
      "— 6-motor asymmetric hexacopter architecture",
      "— Designed around an eVTOL operating mission"]),
    ("②", "Objective",
     ["A single MDO framework that simultaneously couples",
      "structure · propulsion · ACS · powertrain · battery",
      "— supports both single- and multi-objective search"]),
    ("③", "Key Contribution",
     ["High-performance C++ evaluation engine",
      "+ Pagmo2 (CMA-ES / NSGA-II) integration",
      "+ real-time 3D visualization (ImGui / Vulkan)"]),
]
y = 0.9
for mark, title_, lines in items:
    card(s, 0.36, y, 4.55, 1.32, CYAN)
    txt(s, 0.5, y + 0.08, 0.5, 0.4, mark, size=18, color=CYAN, bold=True)
    txt(s, 1.02, y + 0.08, 3.7, 0.3, title_, size=13, color=NAVY, bold=True)
    txt(s, 1.02, y + 0.42, 3.78, 0.85,
        [[(l, {})] for l in lines], size=10, color=GRAY, sp_after=1, line=1.0)
    y += 1.46

# right: stat cards
txt(s, 5.18, 0.82, 4.4, 0.3, "Framework at a Glance", size=14, color=NAVY, bold=True)
stats = [("5", "analysis disciplines", "structure · ACS · powertrain · battery · packaging", CYAN),
         ("16", "active design variables", "geometry 4 + structure 2 + battery 1 + positions 9", ORANGE),
         ("18", "hard constraints", "geometry · packaging · ACS · structural · battery", GREEN),
         ("2", "optimization algorithms", "SOO (CMA-ES)  +  MOO (NSGA-II)", BLUE)]
y = 1.22
for big, lab, sub, c in stats:
    card(s, 5.18, y, 4.46, 0.92, c, fill=FILL1)
    txt(s, 5.34, y, 1.0, 0.92, big, size=34, color=c, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, 6.4, y + 0.14, 3.1, 0.3, lab, size=12, color=NAVY, bold=True)
    txt(s, 6.4, y + 0.46, 3.18, 0.4, sub, size=8.5, color=GRAY, line=0.95)
    y += 1.0

# =============================================================================
# SLIDE 3 — Design Variables
# =============================================================================
s = slide(); header(s, "Optimization Design Variables (16 active)", "03"); footer(s)


ROW_PITCH = 0.46


def var_group(l, t, w, accent, title_, rows):
    h = 0.30 + ROW_PITCH * len(rows)
    card(s, l, t, w, h, accent)
    txt(s, l + 0.16, t + 0.05, w - 0.2, 0.26, title_, size=12, color=NAVY, bold=True)
    yy = t + 0.34
    for name, info in rows:
        txt(s, l + 0.16, yy, w - 0.24, 0.22, name, size=10.5, color=accent, bold=True, font="Consolas")
        txt(s, l + 0.16, yy + 0.205, w - 0.24, 0.2, info, size=8.5, color=GRAY)
        yy += ROW_PITCH
    return h


COL1, COL2, COL3 = 0.36, 3.5, 6.64
ROW1, ROW2 = 0.72, 2.92
var_group(COL1, ROW1, 3.0, CYAN, "Geometry", [
    ("Lx", "default 2.65 m  ·  range [2.0 – 5.0]"),
    ("Lyi", "default 2.65 m  ·  range [2.0 – 5.0]"),
    ("Lyo", "default 5.50 m  ·  range [2.5 – 9.0]"),
    ("T_max", "default 12000 N  ·  range [8k – 20k]"),
])
var_group(COL2, ROW1, 3.0, ORANGE, "Structure & Energy", [
    ("arm_outer_radius", "default 0.08 m  ·  [0.02 – 0.15]"),
    ("arm_wall_thickness", "default 0.005 m  ·  [0.001 – 0.020]"),
    ("m_bat", "default 400 kg  ·  range [100 – 1000]"),
])
var_group(COL3, ROW1, 3.0, GREEN, "Passenger Position", [
    ("pax_x", "0.00 m  ·  range [-0.25, +0.25]"),
    ("pax_y", "0.00 m  ·  range [-0.42, +0.42]"),
    ("pax_z", "-0.10 m  ·  range [-0.50, +0.40]"),
])
var_group(COL2, ROW2, 3.0, BLUE, "Cargo Position", [
    ("cargo_x", "0.00 m  ·  range [-0.50, +0.50]"),
    ("cargo_y", "0.00 m  ·  range [-0.45, +0.45]"),
    ("cargo_z", "0.75 m  ·  range [0.30, 1.00]"),
])
var_group(COL3, ROW2, 3.0, BLUE_LT, "Battery Position", [
    ("bat_x", "0.00 m  ·  range [-0.50, +0.50]"),
    ("bat_y", "0.00 m  ·  range [-0.55, +0.55]"),
    ("bat_z", "-0.90 m  ·  range [-1.00, +0.85]"),
])
# Variable-handling card (normalization + placement semantics)
card(s, COL1, ROW2, 3.0, 1.66, GRAY_LT, fill=FILL2)
txt(s, COL1 + 0.16, ROW2 + 0.05, 2.7, 0.26, "Variable handling", size=11.5, color=NAVY, bold=True)
txt(s, COL1 + 0.16, ROW2 + 0.36, 2.72, 1.25, [
    [("• All 16 variables normalized to ", {"size": 8.5, "color": GRAY}),
     ("[0,1]", {"size": 8.5, "color": NAVY, "bold": True}),
     ("; ", {"size": 8.5, "color": GRAY}),
     ("DesignVectorMapper", {"size": 8.5, "color": NAVY, "bold": True}),
     (" → physical units.", {"size": 8.5, "color": GRAY})],
    [("• pax / cargo / battery are ", {"size": 8.5, "color": GRAY}),
     ("element-owned body-frame placement variables", {"size": 8.5, "color": ORANGE, "bold": True}),
     (" (each element encodes them in its local_pose_) — this is why packaging constraints and the 3D viewer matter.",
      {"size": 8.5, "color": GRAY})],
], line=1.04, sp_after=4)
txt(s, 0.36, 4.78, 9.3, 0.3,
    [[("★  Fixed (inactive) parameters: ", {"size": 9, "color": ORANGE, "bold": True}),
      ("cT = 0.03,  d_prop = 0.40 m,  m_pax = 600 kg,  m_cargo = 200 kg,  m_instrument = 50 kg  — mission requirements, not optimized.",
       {"size": 9, "color": GRAY})]])

# =============================================================================
# SLIDE 4 — Objectives & Hard Constraints
# =============================================================================
s = slide(); header(s, "Objective Functions & Hard Constraints", "04"); footer(s)
txt(s, 0.36, 0.72, 4.4, 0.3, "Objectives  (weighted-sum SOO / Pareto MOO)", size=13, color=NAVY, bold=True)
objs = [("mass", "0.20", "total vehicle mass (normalized)"),
        ("power", "0.20", "hover power proxy"),
        ("fault_thrust", "0.25", "single-motor-fault thrust-margin penalty"),
        ("fault_alloc", "0.25", "control-allocation degradation ratio"),
        ("hover_nom", "0.10", "nominal hover thrust utilization"),
        ("cabin_space_penalty", "0.10", "cabin packing / space penalty"),
        ("acs_margin_penalty", "0.10", "ACS directional-margin penalty (≥0 healthy)")]
y = 1.12
for name, w, desc in objs:
    card(s, 0.36, y, 4.5, 0.5, CYAN, fill=FILL1)
    txt(s, 0.5, y + 0.04, 2.3, 0.24, name, size=10.5, color=NAVY, bold=True, font="Consolas")
    rect(s, 2.9, y + 0.07, 0.92, 0.3, fill=ORANGE)
    txt(s, 2.9, y + 0.07, 0.92, 0.3, "w = " + w, size=9, color=WHITE, bold=True,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, 0.5, y + 0.27, 4.3, 0.22, desc, size=8.5, color=GRAY)
    y += 0.585
txt(s, 0.36, 4.84, 4.5, 0.4,
    [[("structural · packaging · structural_safety carry weight 0 → analysis-only metrics.",
       {"size": 8.5, "color": GRAY_LT})]])

# constraints
txt(s, 5.05, 0.72, 4.6, 0.3, "Hard Constraints  (large penalty on violation)", size=13, color=NAVY, bold=True)
cons = [
    ("GEO", "parameter_bounds  ≤ 0", CYAN),
    ("GEO", "minimum_geometry_margin  ≥ 0", CYAN),
    ("GEO", "rotor_clearance  ≥ 0", CYAN),
    ("PACK", "payload / battery / occupant in cabin", BLUE),
    ("PACK", "battery-payload & component non-overlap", BLUE),
    ("CG", "cg_envelope  |x|≤0.40, |y|≤0.25 m", BLUE_LT),
    ("ACS", "failed_hover_gamma  ≥ 1.5", ORANGE),
    ("ACS", "fault_allocation_ratio  ≥ 0.05", ORANGE),
    ("ACS", "all_faults_hover_feasible  ≥ 0", ORANGE),
    ("ACS", "fault_directional_margin  ≥ 0", ORANGE),
    ("STR", "arm_yield_failure  ≥ 1.5", GREEN),
    ("STR", "arm_tip_deflection  ≤ 0.10 m", GREEN),
    ("STR", "arm_tip_rotation  ≤ 0.10 rad", GREEN),
    ("BAT", "battery_energy_reserve  ≥ 0", BLUE),
    ("BAT", "battery_crate_limit  ≤ 5 /h", BLUE),
]
y = 1.12
for tag, text, c in cons:
    rect(s, 5.05, y, 0.62, 0.235, fill=c)
    txt(s, 5.05, y, 0.62, 0.235, tag, size=8, color=WHITE, bold=True,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, 5.78, y - 0.005, 3.9, 0.24, text, size=9.5, color=INK, font="Consolas",
        anchor=MSO_ANCHOR.MIDDLE)
    y += 0.252
txt(s, 5.05, 4.92, 4.6, 0.3,
    [[("18 architecture-level hard constraints total ", {"size": 8.5, "color": GRAY, "bold": True}),
      ("(packaging row groups 5 sub-checks).", {"size": 8.5, "color": GRAY_LT})]])

# =============================================================================
# SLIDE 5 — Evaluation Pipeline
# =============================================================================
s = slide(); header(s, "Multidisciplinary Evaluation Pipeline (Stage1Evaluator)", "05"); footer(s)
steps = [
    ("Step 01", "Vehicle Scaling Model", ["Mass · CG · inertia · thrust", "battery-mass feedback included"], CYAN),
    ("Step 02", "ACS Analyzer", ["Nominal + 6 fault hover trims (LP)", "margins · retention · feasibility"], ORANGE),
    ("Step 03", "Structural Network", ["Multi-member network, 8 load cases", "von Mises SF · tip deflection"], GREEN),
    ("Step 04", "Powertrain Evaluator", ["Actuator-disk power model", "per-motor electrical power [W]"], BLUE),
    ("Step 05", "Battery Evaluator", ["Energy budget · reserve fraction", "C-rate (mission-aware)"], BLUE_LT),
    ("Step 06", "Constraints + Objective", ["18 hard-constraint checks", "weighted-sum aggregation"], NAVY),
]
positions = [(0.36, 0.85), (3.58, 0.85), (6.8, 0.85), (0.36, 2.55), (3.58, 2.55), (6.8, 2.55)]
for (l, t), (st, name, lines, c) in zip(positions, steps):
    card(s, l, t, 2.84, 1.5, c)
    rect(s, l + 0.16, t + 0.12, 0.78, 0.3, fill=c)
    txt(s, l + 0.16, t + 0.12, 0.78, 0.3, st, size=10, color=WHITE, bold=True,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, l + 1.02, t + 0.12, 1.7, 0.32, name, size=11.5, color=NAVY, bold=True,
        anchor=MSO_ANCHOR.MIDDLE)
    txt(s, l + 0.18, t + 0.56, 2.6, 0.85, [[(x, {})] for x in lines], size=9.5, color=GRAY,
        line=1.05, sp_after=1)
# call-order note
txt(s, 0.36, 4.18, 9.3, 0.3,
    [[("Call order matters: ", {"size": 9.5, "color": ORANGE, "bold": True}),
      ("ACS runs before structural (fault trims become load cases); packaging runs after structural.",
       {"size": 9.5, "color": GRAY})]])
rect(s, 0.7, 4.6, 8.6, 0.46, fill=NAVY)
txt(s, 0.7, 4.6, 8.6, 0.46,
    [[("▶  EvaluationResult  =  PhysicalModel + AcsResult + PowertrainResult + BatteryResult + constraint_results + combined_objective",
       {"size": 9.5, "color": WHITE, "bold": True})]],
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# =============================================================================
# SLIDE 6 — Phase 1 Structural & ACS
# =============================================================================
s = slide(); header(s, "Phase 1 — Structural Analysis & ACS", "06"); footer(s)
card(s, 0.36, 0.78, 4.5, 2.1, GREEN)
txt(s, 0.54, 0.86, 4.2, 0.3, "Structural Network Analyzer", size=13, color=NAVY, bold=True)
txt(s, 0.54, 1.22, 4.25, 1.6, [
    [("▸ Hollow circular tube cross-section per arm", {})],
    [("▸ Outer radius r₀, wall t  →  second moment I", {})],
    [("▸ Max bending load = motor thrust × arm length", {})],
    [("▸ 8 load cases: max-thrust, hover, fault_0–5", {})],
    [("▸ Safety factor SF = σ_yield / σ_bending ≥ 1.5", {})],
    [("▸ Euler–Bernoulli tip deflection & rotation", {})],
], size=10, color=GRAY, line=1.1, sp_after=2)
card(s, 5.0, 0.78, 4.64, 2.1, ORANGE)
txt(s, 5.18, 0.86, 4.3, 0.3, "Attainable Control Set Analyzer", size=13, color=NAVY, bold=True)
txt(s, 5.18, 1.22, 4.35, 1.6, [
    [("▸ Allocation matrix B ∈ ℝ⁴ˣ⁶ (thrust·roll·pitch·yaw)", {})],
    [("▸ LP hover trim: nominal + 6 single-motor faults", {})],
    [("▸ Directional margin  m(d) = h_U(d) − dᵀ·u_req", {})],
    [("▸ PFWAR — probability-weighted ACS retention", {})],
    [("▸ WCFR — worst-case fault retention", {})],
    [("▸ FII — fault isotropy index (FII < 1 is ideal)", {})],
], size=10, color=GRAY, line=1.1, sp_after=2)
# invariant box
rect(s, 0.36, 3.02, 9.28, 0.78, fill=NAVY)
rect(s, 0.36, 3.02, 0.1, 0.78, fill=CYAN)
txt(s, 0.6, 3.08, 9.0, 0.3, "ACS correctness invariant", size=11, color=CYAN2, bold=True)
txt(s, 0.6, 3.4, 9.0, 0.36,
    [[("margin(d, u) = h_U(d) − dᵀ·u  ≥ 0  feasible      h_U(d) = Σⱼ max(0, dᵀ·bⱼ)·f_maxⱼ      ",
       {"size": 10, "color": WHITE, "font": "Consolas"}),
      ("(never use the sign-flipped form).", {"size": 9.5, "color": GRAY_LT})]])
# output files
txt(s, 0.36, 3.96, 5.0, 0.26, "ACS output files  (--plot-acs)", size=11, color=NAVY, bold=True)
files = [("acs_fig5_hover_poly.svg", "hover moment polygon"),
         ("acs_fig6_margins.svg", "per-direction margin bars"),
         ("acs_retention.svg", "per-motor ACS retention"),
         ("acs_vertices.csv", "4D vertex cloud")]
y = 4.26
for i, (f, d) in enumerate(files):
    col = 0.36 + (i % 2) * 4.7
    row = y + (i // 2) * 0.34
    txt(s, col, row, 2.4, 0.24, f, size=9, color=CYAN, bold=True, font="Consolas")
    txt(s, col + 2.4, row, 2.2, 0.24, "— " + d, size=9, color=GRAY)

# =============================================================================
# SLIDE 7 — Phase 2 Powertrain & Battery
# =============================================================================
s = slide(); header(s, "Phase 2 — Powertrain & Battery Energy Model", "07"); footer(s)
card(s, 0.36, 0.78, 4.5, 3.05, BLUE)
txt(s, 0.54, 0.86, 4.2, 0.3, "Powertrain  (PowertrainEvaluator)", size=13, color=NAVY, bold=True)
pw = [("Effective disk radius", "r_eff = max_arm_length / 2"),
      ("Electrical power (rotor i)", "Pᵢ = Tᵢ^1.5 / (η_tot · √(2ρ·A_eff))"),
      ("Total nominal power", "P_nom = Σ Pᵢ  (nominal hover trim)"),
      ("Total fault power", "P_fault = Σ Pᵢ  (worst fault trim)")]
yy = 1.24
for lab, eq in pw:
    txt(s, 0.54, yy, 4.2, 0.22, lab, size=9, color=GRAY)
    txt(s, 0.54, yy + 0.2, 4.2, 0.24, eq, size=10.5, color=NAVY, bold=True, font="Consolas")
    yy += 0.56
txt(s, 0.54, 3.5, 4.25, 0.3,
    [[("η_tot = FoM·η_motor·η_esc = 0.65·0.85·0.95.  d_prop drives yaw-torque scaling only.",
       {"size": 8.5, "color": GRAY})]])

card(s, 5.0, 0.78, 4.64, 3.05, ORANGE)
txt(s, 5.18, 0.86, 4.3, 0.3, "Battery  (BatteryEvaluator)", size=13, color=NAVY, bold=True)
bt = [("Available energy", "E_avail = η_pack · DoD · m_bat · e_spec"),
      ("Mission energy demand", "E_req = P_nom·t_nom + P_fault·t_emg"),
      ("Energy reserve", "reserve = (E_avail − E_req)/E_avail ≥ 0"),
      ("C-rate", "c_rate = P_peak / E_avail  ≤  5.0 /h")]
yy = 1.24
for lab, eq in bt:
    txt(s, 5.18, yy, 4.3, 0.22, lab, size=9, color=GRAY)
    txt(s, 5.18, yy + 0.2, 4.3, 0.24, eq, size=10.5, color=NAVY, bold=True, font="Consolas")
    yy += 0.56
txt(s, 5.18, 3.5, 4.3, 0.3,
    [[("e_spec = 250 Wh/kg (NMC811) · DoD = 0.85 · η_pack = 0.95 · t_nom = 30 min · t_emg = 1 min",
       {"size": 8.5, "color": GRAY})]])

# feedback loop
rect(s, 0.36, 4.0, 9.28, 1.05, fill=FILL2)
rect(s, 0.36, 4.0, 0.1, 1.05, fill=CYAN)
txt(s, 0.56, 4.08, 9.0, 0.28, "Battery-mass feedback loop", size=12, color=NAVY, bold=True)
txt(s, 0.56, 4.4, 9.0, 0.3,
    [[("m_bat ↑  →  total mass ↑  →  hover thrust ↑  →  power ↑  →  energy demand ↑  →  battery constraints re-checked",
       {"size": 10, "color": INK, "bold": True})]])
txt(s, 0.56, 4.72, 9.0, 0.3,
    [[("Below the minimum viable battery mass the energy reserve goes negative — the optimizer finds the nonlinear balance point.",
       {"size": 9, "color": GRAY})]])

# =============================================================================
# SLIDE 8 — Optimization Algorithms & Loop
# =============================================================================
s = slide(); header(s, "Optimization — SOO / MOO and the Search Loop", "08"); footer(s)
card(s, 0.36, 0.78, 4.5, 3.0, CYAN)
txt(s, 0.54, 0.86, 4.2, 0.3, "SOO — Single-Objective (SooRunner)", size=12.5, color=NAVY, bold=True)
soo = [("Algorithm", "CMA-ES (Pagmo2)"), ("Population", "24 (default)"),
       ("Generations", "40 (default)"), ("Search space", "normalized [0,1]¹⁶"),
       ("Feasibility", "all hard constraints + ACS LP"),
       ("Output", "baseline / best_raw / best_feasible")]
yy = 1.2
for k, v in soo:
    txt(s, 0.54, yy, 1.55, 0.26, k, size=9.5, color=GRAY, bold=True)
    txt(s, 2.12, yy, 2.7, 0.26, v, size=9.5, color=NAVY)
    yy += 0.30
txt(s, 0.54, 3.08, 4.25, 0.6, [
    [("Penalty strategy:", {"size": 9, "color": ORANGE, "bold": True})],
    [("infeasible → 1e6 + Σ(weightᵢ·violationᵢ) + 1e4", {"size": 8.5, "color": GRAY, "font": "Consolas"})],
    [("feasible → combined_objective (weighted sum)", {"size": 8.5, "color": GRAY, "font": "Consolas"})],
], line=1.05, sp_after=1)

card(s, 5.0, 0.78, 4.64, 3.0, ORANGE)
txt(s, 5.18, 0.86, 4.3, 0.3, "MOO — Multi-Objective (MooRunner)", size=12.5, color=NAVY, bold=True)
moo = [("Algorithm", "NSGA-II (Pagmo2)"), ("Population", "48 (default)"),
       ("Generations", "60 (default)"), ("Objectives", "mass, power, fault_alloc"),
       ("Output", "Pareto front (feasible subset)"),
       ("Knee analysis", "ParetoAnalyzer (knee point)")]
yy = 1.2
for k, v in moo:
    txt(s, 5.18, yy, 1.7, 0.26, k, size=9.5, color=GRAY, bold=True)
    txt(s, 6.9, yy, 2.7, 0.26, v, size=9.5, color=NAVY)
    yy += 0.30
txt(s, 5.18, 3.08, 4.3, 0.6, [
    [("Trade-off surface (3 objectives):", {"size": 9, "color": BLUE, "bold": True})],
    [("min mass  ←→  min power  ←→  max fault tolerance", {"size": 8.5, "color": GRAY})],
    [("designer selects the operating point directly", {"size": 8.5, "color": GRAY})],
], line=1.05, sp_after=1)

# loop chain
boxes = ["x̂ ∈ [0,1]¹⁶", "DesignVectorMapper", "HexacopterArchitecture", "Stage1Evaluator", "Pagmo fitness", "next generation"]
bw = 1.46; gap = 0.12; x = 0.36
for i, b in enumerate(boxes):
    c = NAVY if i in (0, 5) else BLUE
    rect(s, x, 4.05, bw, 0.5, fill=FILL1, line=c, lw=1.0)
    txt(s, x, 4.05, bw, 0.5, b, size=9, color=NAVY, bold=True, align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE, font="Consolas")
    if i < len(boxes) - 1:
        txt(s, x + bw - 0.02, 4.05, 0.16, 0.5, "→", size=12, color=ORANGE, bold=True,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    x += bw + gap
txt(s, 0.36, 4.66, 9.3, 0.3,
    [[("Iterative loop: search in normalized space → map to physical units → evaluate all disciplines → feed fitness back.",
       {"size": 9, "color": GRAY})]], align=PP_ALIGN.CENTER)

# =============================================================================
# SLIDE 9 — Code Architecture
# =============================================================================
s = slide(); header(s, "C++ Code Architecture (FrameworkCpp)", "09"); footer(s)
card(s, 0.36, 0.78, 4.2, 4.1, CYAN, fill=FILL1)
txt(s, 0.54, 0.86, 3.9, 0.3, "Directory layout", size=12.5, color=NAVY, bold=True)
tree = ["FrameworkCpp/", "  ├─ app/            (main.cpp · CLI modes)",
        "  ├─ include/ + src/",
        "  │   ├─ core/        Architecture · Elements",
        "  │   ├─ physics/     Structural · ACS · Power · Battery",
        "  │   ├─ evaluation/  Stage1Evaluator",
        "  │   ├─ optimization/ SOO · MOO · Pagmo",
        "  │   ├─ analysis/    CSV · JSON · Pareto · ACS plot",
        "  │   ├─ mission/     MissionProfile · Evaluator",
        "  │   └─ calibration/ Calibrator (Nelder–Mead)",
        "  ├─ visualization/   ImGui / Vulkan 3D viewer",
        "  └─ docs/            CODEBASE.md · Action.md …"]
txt(s, 0.54, 1.2, 3.95, 3.6, [[(t, {})] for t in tree], size=8.5, color=GRAY,
    font="Consolas", line=1.12, sp_after=1)

txt(s, 4.78, 0.82, 5.0, 0.3, "Core classes", size=12.5, color=NAVY, bold=True)
cls = [("HexacopterArchitecture", "root: parameter · element · constraint registries"),
       ("DesignParameter", "design variable (value · bounds · norm · id)"),
       ("Stage1Evaluator", "runs the 6-step multidisciplinary pipeline"),
       ("PagmoProblemAdapter", "Pagmo2 adapter (shared by SOO / MOO)"),
       ("SooRunner / MooRunner", "CMA-ES / NSGA-II execution wrappers"),
       ("ConstraintRegistry", "hard-constraint evaluation + penalties"),
       ("CsvExporter", "results to CSV / JSON in physical units"),
       ("AcsPlotter", "SVG ACS visualization (no dependencies)"),
       ("ArchitectureViewerApp", "ImGui / Vulkan real-time 3D viewer")]
y = 1.16
for name, desc in cls:
    rect(s, 4.78, y, 4.86, 0.385, fill=FILL1)
    rect(s, 4.78, y, 0.06, 0.385, fill=BLUE)
    txt(s, 4.92, y + 0.03, 4.7, 0.22, name, size=10, color=NAVY, bold=True, font="Consolas")
    txt(s, 4.92, y + 0.22, 4.7, 0.18, desc, size=8, color=GRAY)
    y += 0.405

# =============================================================================
# SLIDE 10 — Execution & Outputs
# =============================================================================
s = slide(); header(s, "How to Run & Output Artifacts", "10"); footer(s)
card(s, 0.36, 0.78, 5.2, 3.0, NAVY, fill=NAVY)
txt(s, 0.54, 0.84, 5.0, 0.3, "CLI commands", size=12.5, color=CYAN2, bold=True)
cmds = [("# baseline design evaluation", "FrameworkCpp.exe eval"),
        ("# single-objective optimization", "FrameworkCpp.exe soo --soo-pop 24 --soo-gen 40"),
        ("# multi-objective optimization", "FrameworkCpp.exe moo --moo-pop 48 --moo-gen 60"),
        ("# full SOO + MOO comparison", "FrameworkCpp.exe compare --output-dir ./out"),
        ("# SOO with live 3D visualization", "FrameworkCpp.exe soo --visualize"),
        ("# mission / calibration subsystems", "FrameworkCpp.exe mission data\\...json")]
yy = 1.18
for c, cmd in cmds:
    txt(s, 0.54, yy, 5.0, 0.18, c, size=8, color=GRAY_LT, font="Consolas")
    txt(s, 0.54, yy + 0.16, 5.0, 0.22, cmd, size=9.5, color=WHITE, font="Consolas")
    yy += 0.42

txt(s, 5.78, 0.82, 3.9, 0.3, "Output files", size=12.5, color=NAVY, bold=True)
outs = [("comparison.csv", "baseline vs optimum comparison"),
        ("soo_parameters.csv", "SOO optimal variables (physical units)"),
        ("pareto_front.csv", "Pareto objectives + safety factors"),
        ("soo_run.json", "full SOO result (with physical model)"),
        ("moo_run.json", "MOO result (Pareto population)"),
        ("acs/*.svg", "ACS margin & retention charts")]
y = 1.16
for f, d in outs:
    rect(s, 5.78, y, 3.86, 0.385, fill=FILL1)
    rect(s, 5.78, y, 0.06, 0.385, fill=ORANGE)
    txt(s, 5.92, y + 0.03, 3.7, 0.22, f, size=9.5, color=NAVY, bold=True, font="Consolas")
    txt(s, 5.92, y + 0.22, 3.7, 0.18, d, size=8, color=GRAY)
    y += 0.43

rect(s, 0.36, 3.96, 9.28, 1.05, fill=FILL2)
rect(s, 0.36, 3.96, 0.1, 1.05, fill=CYAN)
txt(s, 0.56, 4.04, 9.0, 0.28, "Real-time 3D visualization  (ArchitectureViewerApp — ImGui + Vulkan)",
    size=11.5, color=NAVY, bold=True)
viz = ["• live arm / motor / rotor / payload rendering", "• updates as the optimization progresses",
       "• loads SOO/MOO results from JSON", "• left panel: per-element mass & position"]
for i, v in enumerate(viz):
    col = 0.56 + (i % 2) * 4.7
    row = 4.36 + (i // 2) * 0.3
    txt(s, col, row, 4.6, 0.26, v, size=9.5, color=GRAY)

# =============================================================================
# SLIDE 11 — Post-Optimization Analysis
# =============================================================================
s = slide(); header(s, "Post-Optimization Analysis & Interpretation", "11"); footer(s)
txt(s, 0.36, 0.70, 9.3, 0.28,
    [[("Each run tracks ", {"size": 10, "color": GRAY}),
      ("raw-best and feasible-best separately", {"size": 10, "color": NAVY, "bold": True}),
      (", then exports artifacts for the designer to interpret the trade space.", {"size": 10, "color": GRAY})]])

# SOO interpretation card
card(s, 0.36, 1.12, 4.5, 1.95, CYAN)
txt(s, 0.54, 1.2, 4.2, 0.28, "SOO outcome  —  three tracked solutions", size=12, color=NAVY, bold=True)
soo_out = [("baseline", "default design vector (reference point)", GRAY_LT),
           ("best_raw", "lowest objective ignoring feasibility", ORANGE),
           ("best_feasible", "lowest objective with all 18 constraints met", GREEN)]
yy = 1.55
for name, desc, c in soo_out:
    rect(s, 0.54, yy + 0.03, 0.14, 0.34, fill=c)
    txt(s, 0.78, yy, 2.0, 0.24, name, size=10.5, color=NAVY, bold=True, font="Consolas")
    txt(s, 0.78, yy + 0.21, 3.95, 0.2, desc, size=8.5, color=GRAY)
    yy += 0.47

# MOO interpretation card
card(s, 5.0, 1.12, 4.64, 1.95, ORANGE)
txt(s, 5.18, 1.2, 4.3, 0.28, "MOO outcome  —  Pareto trade-off", size=12, color=NAVY, bold=True)
txt(s, 5.18, 1.55, 4.35, 1.4, [
    [("• ", {"size": 9.5, "color": ORANGE, "bold": True}),
     ("Feasible Pareto subset", {"size": 9.5, "color": NAVY, "bold": True}),
     (" — non-dominated designs that satisfy every hard constraint.", {"size": 9.5, "color": GRAY})],
    [("• ", {"size": 9.5, "color": ORANGE, "bold": True}),
     ("Knee-point selection", {"size": 9.5, "color": NAVY, "bold": True}),
     (" via ParetoAnalyzer — best mass / power / fault-tolerance compromise.", {"size": 9.5, "color": GRAY})],
    [("• ", {"size": 9.5, "color": ORANGE, "bold": True}),
     ("Designer picks the operating point along the front directly.", {"size": 9.5, "color": GRAY})],
], line=1.08, sp_after=4)

# Post-analysis artifact flow
rect(s, 0.36, 3.28, 9.28, 1.62, fill=FILL2)
rect(s, 0.36, 3.28, 0.1, 1.62, fill=CYAN)
txt(s, 0.56, 3.36, 9.0, 0.28, "Post-analysis interpretation flow", size=12, color=NAVY, bold=True)
arts = [("ACS SVGs", "directional margins,\nfault retention", CYAN),
        ("comparison.csv", "baseline vs optimum,\nside by side", BLUE),
        ("*_parameters.csv", "winning design vector\nin physical units", ORANGE),
        ("soo/moo_run.json", "full result +\nphysical model", GREEN),
        ("3D viewer replay", "load JSON, inspect\nplacement & masses", BLUE_LT)]
bw, gap = 1.74, 0.12
x = 0.56
for i, (name, desc, c) in enumerate(arts):
    rect(s, x, 3.72, bw, 1.0, fill=WHITE, line=c, lw=1.25)
    rect(s, x, 3.72, bw, 0.06, fill=c)
    txt(s, x + 0.08, 3.84, bw - 0.16, 0.4, name, size=9.5, color=NAVY, bold=True,
        align=PP_ALIGN.CENTER, font="Consolas")
    txt(s, x + 0.08, 4.18, bw - 0.16, 0.5,
        [[(ln, {})] for ln in desc.split("\n")], size=8, color=GRAY, align=PP_ALIGN.CENTER, line=1.0)
    if i < len(arts) - 1:
        txt(s, x + bw - 0.02, 3.72, 0.16, 1.0, "→", size=13, color=ORANGE, bold=True,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    x += bw + gap
txt(s, 0.36, 4.98, 9.3, 0.3,
    [[("Generate everything at once:  ", {"size": 9, "color": GRAY}),
      ("FrameworkCpp.exe compare --output-dir ./out", {"size": 9, "color": NAVY, "bold": True, "font": "Consolas"}),
      ("   (runs SOO + MOO and writes all artifacts).", {"size": 9, "color": GRAY})]])

# =============================================================================
# SLIDE 12 — Conclusion
# =============================================================================
s = slide()
rect(s, 0, 0, 10, 5.625, fill=NAVY)
rect(s, 0, 0, 10, 0.10, fill=ORANGE)
txt(s, 0.5, 0.45, 9.0, 0.7, "Conclusion & Future Work", size=30, color=WHITE, bold=True, font=FONT_L)
rect(s, 0.52, 1.2, 2.4, 0.04, fill=CYAN)
checks = [
    "C++ high-performance MDO framework — 5 coupled disciplines (structure · ACS · propulsion · battery · packaging)",
    "16 design variables, 18 hard constraints → Pagmo2 SOO (CMA-ES) + MOO (NSGA-II)",
    "Phase 2: battery energy / C-rate constraints and powertrain model integrated",
    "3-DOF placement of passenger · cargo · battery with AABB cabin-packaging constraints",
    "ImGui / Vulkan real-time 3D viewer + automatic ACS SVG plots",
    "Mission and calibration subsystems (multi-segment energy, flight-log parameter ID)",
]
y = 1.45
for c in checks:
    rect(s, 0.5, y + 0.02, 0.26, 0.26, fill=CYAN, shape=MSO_SHAPE.OVAL)
    txt(s, 0.5, y + 0.02, 0.26, 0.26, "✓", size=11, color=NAVY, bold=True,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, 0.9, y, 8.7, 0.34, c, size=11.5, color=FILL2, anchor=MSO_ANCHOR.MIDDLE)
    y += 0.45
txt(s, 0.5, 4.3, 2.0, 0.28, "Future work", size=13, color=ORANGE, bold=True)
fut = ["→ higher-fidelity CFD / FEM coupling", "→ range and noise objectives",
       "→ flight-data-driven model calibration", "→ multi-mission (cruise + hover) optimization"]
for i, f in enumerate(fut):
    col = 0.5 + (i % 2) * 4.7
    row = 4.62 + (i // 2) * 0.3
    txt(s, col, row, 4.6, 0.26, f, size=10.5, color=FILL2)
rect(s, 0, 5.525, 10, 0.10, fill=CYAN)

out = "UAM_MDO_Presentation_EN.pptx"
prs.save(out)
print("saved:", out, "| slides:", len(prs.slides._sldIdLst))
